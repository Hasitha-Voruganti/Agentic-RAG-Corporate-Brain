"""
ingestion/parser.py — Multi-format document parser
Supports: PDF, DOCX, XLSX, PPTX, TXT
With OCR fallback for scanned PDFs
"""
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from dataclasses import dataclass, field
from typing import Optional

import pdfplumber
import pytesseract
from pdf2image import convert_from_path
import docx
import openpyxl
from loguru import logger


@dataclass
class ParsedChunk:
    content: str
    page_num: Optional[int] = None
    chunk_type: str = "text"
    metadata: dict = field(default_factory=dict)


class DocumentParser:
    """
    Parse any supported document format into a list of ParsedChunk objects.
    Supported: .pdf, .docx, .xlsx, .pptx, .txt
    """

    def parse(self, file_path: str) -> list[ParsedChunk]:
        ext = Path(file_path).suffix.lower()
        parsers = {
            ".pdf":  self._parse_pdf,
            ".docx": self._parse_docx,
            ".xlsx": self._parse_xlsx,
            ".pptx": self._parse_pptx,
            ".txt":  self._parse_txt,
        }
        if ext not in parsers:
            raise ValueError(
                f"Unsupported file type: {ext}. "
                f"Supported types: {list(parsers.keys())}"
            )
        logger.info(f"Parsing {ext} file: {Path(file_path).name}")
        return parsers[ext](file_path)

    # ── PDF ────────────────────────────────────────────────────────────────
    def _parse_pdf(self, path: str) -> list[ParsedChunk]:
        chunks = []
        try:
            with pdfplumber.open(path) as pdf:
                total_pages = len(pdf.pages)
                logger.info(f"PDF has {total_pages} pages")

                for page_num, page in enumerate(pdf.pages, start=1):
                    # Extract text
                    text = page.extract_text()
                    if text and text.strip():
                        chunks.append(ParsedChunk(
                            content=text.strip(),
                            page_num=page_num,
                            chunk_type="text",
                            metadata={"source_page": page_num}
                        ))

                    # Extract tables
                    try:
                        tables = page.extract_tables()
                        for table in tables:
                            if not table:
                                continue
                            rows = []
                            for row in table:
                                cleaned = [
                                    str(cell or "").strip()
                                    for cell in row
                                ]
                                row_text = " | ".join(
                                    c for c in cleaned if c
                                )
                                if row_text:
                                    rows.append(row_text)
                            if rows:
                                chunks.append(ParsedChunk(
                                    content="\n".join(rows),
                                    page_num=page_num,
                                    chunk_type="table",
                                    metadata={
                                        "source_page": page_num,
                                        "is_table": True
                                    }
                                ))
                    except Exception as te:
                        logger.warning(f"Table extraction failed page {page_num}: {te}")

                    # OCR fallback for scanned pages
                    if not text or len(text.strip()) < 50:
                        logger.info(f"Page {page_num}: sparse text, running OCR")
                        ocr_text = self._ocr_page(path, page_num)
                        if ocr_text and ocr_text.strip():
                            chunks.append(ParsedChunk(
                                content=ocr_text.strip(),
                                page_num=page_num,
                                chunk_type="image_ocr",
                                metadata={
                                    "source_page": page_num,
                                    "ocr": True
                                }
                            ))

        except Exception as e:
            logger.error(f"PDF parse failed for {path}: {e}")

        logger.info(f"PDF parsed: {len(chunks)} chunks extracted")
        return chunks

    def _ocr_page(self, pdf_path: str, page_num: int) -> str:
        """Run Tesseract OCR on a single PDF page."""
        try:
            images = convert_from_path(
                pdf_path,
                first_page=page_num,
                last_page=page_num,
                dpi=200
            )
            if images:
                return pytesseract.image_to_string(images[0], lang="eng")
        except Exception as e:
            logger.warning(f"OCR failed for page {page_num}: {e}")
        return ""

    # ── DOCX ───────────────────────────────────────────────────────────────
    def _parse_docx(self, path: str) -> list[ParsedChunk]:
        chunks = []
        try:
            doc = docx.Document(path)

            # Extract all paragraph text
            texts = []
            for para in doc.paragraphs:
                try:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
                except Exception:
                    continue

            if texts:
                chunks.append(ParsedChunk(
                    content="\n".join(texts),
                    chunk_type="text"
                ))

            # Extract tables
            for table_idx, table in enumerate(doc.tables):
                try:
                    rows = []
                    seen = set()
                    for row in table.rows:
                        cells = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            cells.append(cell_text)
                        row_text = " | ".join(c for c in cells if c)
                        # Avoid duplicate merged cell rows
                        if row_text and row_text not in seen:
                            rows.append(row_text)
                            seen.add(row_text)
                    if rows:
                        chunks.append(ParsedChunk(
                            content="\n".join(rows),
                            chunk_type="table",
                            metadata={
                                "is_table": True,
                                "table_index": table_idx
                            }
                        ))
                except Exception as te:
                    logger.warning(f"DOCX table {table_idx} extraction failed: {te}")

        except Exception as e:
            logger.warning(f"DOCX primary parse failed: {e} — trying XML fallback")
            chunks = self._parse_docx_fallback(path)

        logger.info(f"DOCX parsed: {len(chunks)} chunks extracted")
        return chunks

    def _parse_docx_fallback(self, path: str) -> list[ParsedChunk]:
        """Extract raw text from docx XML as fallback."""
        chunks = []
        try:
            with zipfile.ZipFile(path) as z:
                with z.open("word/document.xml") as f:
                    tree = ET.parse(f)
                    root = tree.getroot()
                    ns = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
                    texts = []
                    for elem in root.iter(f"{ns}t"):
                        if elem.text and elem.text.strip():
                            texts.append(elem.text.strip())
                    content = " ".join(texts)
                    if content.strip():
                        chunks.append(ParsedChunk(
                            content=content.strip(),
                            chunk_type="text",
                            metadata={"fallback": True}
                        ))
            logger.info("DOCX XML fallback succeeded")
        except Exception as e:
            logger.error(f"DOCX XML fallback failed: {e}")
        return chunks

    # ── XLSX ───────────────────────────────────────────────────────────────
    def _parse_xlsx(self, path: str) -> list[ParsedChunk]:
        chunks = []
        try:
            wb = openpyxl.load_workbook(
                path, read_only=True, data_only=True
            )
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                header = None

                for i, row in enumerate(ws.iter_rows(values_only=True)):
                    row_values = [
                        str(v).strip() if v is not None else ""
                        for v in row
                    ]
                    if not any(row_values):
                        continue

                    if i == 0:
                        header = row_values
                        rows.append("COLUMNS: " + " | ".join(header))
                    else:
                        if header:
                            row_text = ", ".join(
                                f"{header[j]}: {v}"
                                for j, v in enumerate(row_values)
                                if v and j < len(header)
                            )
                        else:
                            row_text = " | ".join(
                                v for v in row_values if v
                            )
                        if row_text:
                            rows.append(row_text)

                if rows:
                    chunks.append(ParsedChunk(
                        content="\n".join(rows),
                        chunk_type="table",
                        metadata={
                            "sheet": sheet_name,
                            "is_table": True
                        }
                    ))
                    logger.info(
                        f"XLSX sheet '{sheet_name}': {len(rows)} rows extracted"
                    )

        except Exception as e:
            logger.error(f"XLSX parse failed: {e}")

        logger.info(f"XLSX parsed: {len(chunks)} chunks extracted")
        return chunks

    # ── PPTX ───────────────────────────────────────────────────────────────
    def _parse_pptx(self, path: str) -> list[ParsedChunk]:
        chunks = []
        try:
            from pptx import Presentation
            from pptx.util import Pt
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            prs = Presentation(path)
            total_slides = len(prs.slides)
            logger.info(f"PPTX has {total_slides} slides")

            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_texts = []
                slide_tables = []

                for shape in slide.shapes:
                    # Extract text from text frames
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            para_text = ""
                            for run in para.runs:
                                if run.text.strip():
                                    para_text += run.text
                            if para_text.strip():
                                slide_texts.append(para_text.strip())

                    # Extract tables from slides
                    if shape.has_table:
                        table_rows = []
                        seen_rows = set()
                        for row in shape.table.rows:
                            row_cells = []
                            for cell in row.cells:
                                cell_text = cell.text.strip()
                                row_cells.append(cell_text)
                            row_text = " | ".join(
                                c for c in row_cells if c
                            )
                            if row_text and row_text not in seen_rows:
                                table_rows.append(row_text)
                                seen_rows.add(row_text)
                        if table_rows:
                            slide_tables.append("\n".join(table_rows))

                    # Extract text from grouped shapes
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        for grouped_shape in shape.shapes:
                            if grouped_shape.has_text_frame:
                                for para in grouped_shape.text_frame.paragraphs:
                                    text = para.text.strip()
                                    if text:
                                        slide_texts.append(text)

                # Add slide text chunk
                if slide_texts:
                    slide_notes = ""
                    try:
                        if (slide.has_notes_slide and
                                slide.notes_slide.notes_text_frame):
                            notes = slide.notes_slide.notes_text_frame.text.strip()
                            if notes:
                                slide_notes = f"\nSpeaker Notes: {notes}"
                    except Exception:
                        pass

                    content = "\n".join(slide_texts) + slide_notes
                    chunks.append(ParsedChunk(
                        content=content,
                        page_num=slide_num,
                        chunk_type="text",
                        metadata={"slide": slide_num}
                    ))

                # Add table chunks from slide
                for table_text in slide_tables:
                    chunks.append(ParsedChunk(
                        content=table_text,
                        page_num=slide_num,
                        chunk_type="table",
                        metadata={
                            "slide": slide_num,
                            "is_table": True
                        }
                    ))

        except ImportError:
            logger.error(
                "python-pptx not installed. "
                "Run: pip install python-pptx"
            )
        except Exception as e:
            logger.error(f"PPTX parse failed: {e}")
            # Fallback: extract via XML
            chunks = self._parse_pptx_fallback(path)

        logger.info(f"PPTX parsed: {len(chunks)} chunks extracted")
        return chunks

    def _parse_pptx_fallback(self, path: str) -> list[ParsedChunk]:
        """Extract raw text from PPTX XML as fallback."""
        chunks = []
        try:
            with zipfile.ZipFile(path) as z:
                slide_files = sorted([
                    name for name in z.namelist()
                    if name.startswith("ppt/slides/slide")
                    and name.endswith(".xml")
                    and "/rels/" not in name
                ])
                for slide_num, slide_file in enumerate(slide_files, start=1):
                    with z.open(slide_file) as f:
                        tree = ET.parse(f)
                        root = tree.getroot()
                        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
                        texts = []
                        for elem in root.iter(f"{{{ns}}}t"):
                            if elem.text and elem.text.strip():
                                texts.append(elem.text.strip())
                        if texts:
                            chunks.append(ParsedChunk(
                                content=" ".join(texts),
                                page_num=slide_num,
                                chunk_type="text",
                                metadata={
                                    "slide": slide_num,
                                    "fallback": True
                                }
                            ))
            logger.info(
                f"PPTX XML fallback: {len(chunks)} slides extracted"
            )
        except Exception as e:
            logger.error(f"PPTX XML fallback failed: {e}")
        return chunks

    # ── TXT ────────────────────────────────────────────────────────────────
    def _parse_txt(self, path: str) -> list[ParsedChunk]:
        chunks = []
        try:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
            if content.strip():
                chunks.append(ParsedChunk(
                    content=content.strip(),
                    chunk_type="text"
                ))
                logger.info("TXT parsed: 1 chunk extracted")
        except Exception as e:
            logger.error(f"TXT parse failed: {e}")
        return chunks


# ── Text Splitter ──────────────────────────────────────────────────────────
class TextSplitter:
    """
    Split large chunks into smaller overlapping chunks
    for better retrieval accuracy.
    """

    def __init__(self, chunk_size: int = 512, overlap: int = 64):
        self.chunk_size = chunk_size
        self.overlap = overlap

    def split(self, chunks: list[ParsedChunk]) -> list[ParsedChunk]:
        result = []
        for chunk in chunks:
            word_count = len(chunk.content.split())
            if word_count <= self.chunk_size:
                result.append(chunk)
            else:
                sub_chunks = self._split_text(chunk.content)
                for i, text in enumerate(sub_chunks):
                    result.append(ParsedChunk(
                        content=text,
                        page_num=chunk.page_num,
                        chunk_type=chunk.chunk_type,
                        metadata={**chunk.metadata, "sub_chunk": i}
                    ))
        return result

    def _split_text(self, text: str) -> list[str]:
        words = text.split()
        chunks = []
        start = 0
        while start < len(words):
            end = min(start + self.chunk_size, len(words))
            chunk = " ".join(words[start:end])
            chunks.append(chunk)
            if end == len(words):
                break
            start += self.chunk_size - self.overlap
        return chunks


# ── Singletons ─────────────────────────────────────────────────────────────
document_parser = DocumentParser()
text_splitter = TextSplitter()